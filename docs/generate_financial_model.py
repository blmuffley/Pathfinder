#!/usr/bin/env python3
"""Generate the comprehensive Pathfinder financial model (XLSX) and investor PPTX.

Produces:
- Pathfinder_Financial_Model.xlsx (10 sheets)
- Pathfinder_Business_Case.pptx (8 slides)

Run: python3 docs/generate_financial_model.py
"""

import os
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side, numbers
from openpyxl.utils import get_column_letter

ICLOUD_BASE = os.path.expanduser(
    "~/Library/Mobile Documents/com~apple~CloudDocs/Projects/Avennorth/Solutions/Pathfinder"
)
ICLOUD_FINANCIAL = os.path.join(ICLOUD_BASE, "Financial")
ICLOUD_PRESENTATIONS = os.path.join(ICLOUD_BASE, "Presentations")
os.makedirs(ICLOUD_FINANCIAL, exist_ok=True)
os.makedirs(ICLOUD_PRESENTATIONS, exist_ok=True)

# ══════════════════════════════════════
# Styles
# ══════════════════════════════════════
HEADER_FILL = PatternFill(start_color="1E1E32", end_color="1E1E32", fill_type="solid")
SECTION_FILL = PatternFill(start_color="E8F5E9", end_color="E8F5E9", fill_type="solid")
INPUT_FILL = PatternFill(start_color="E3F2FD", end_color="E3F2FD", fill_type="solid")
TOTAL_FILL = PatternFill(start_color="FFF3E0", end_color="FFF3E0", fill_type="solid")
HEADER_FONT = Font(name="Calibri", size=10, bold=True, color="FFFFFF")
TITLE_FONT = Font(name="Calibri", size=14, bold=True, color="1C1917")
SECTION_FONT = Font(name="Calibri", size=11, bold=True, color="2E7D32")
DATA_FONT = Font(name="Calibri", size=10, color="333333")
BOLD_FONT = Font(name="Calibri", size=10, bold=True, color="333333")
TOTAL_FONT = Font(name="Calibri", size=10, bold=True, color="E65100")
MONEY_FMT = '#,##0'
MONEY_K = '#,##0,"K"'
MONEY_M = '#,##0.0,,"M"'
PCT_FMT = '0.0%'
RATIO_FMT = '0.0"x"'
MONTHS_FMT = '0.0'
thin_border = Border(
    left=Side(style='thin', color='CCCCCC'), right=Side(style='thin', color='CCCCCC'),
    top=Side(style='thin', color='CCCCCC'), bottom=Side(style='thin', color='CCCCCC')
)

def hdr(ws, row, col, val):
    c = ws.cell(row=row, column=col, value=val)
    c.font = HEADER_FONT; c.fill = HEADER_FILL; c.alignment = Alignment(horizontal='center'); c.border = thin_border
    return c

def cell(ws, row, col, val, fmt=None, font=DATA_FONT, fill=None):
    c = ws.cell(row=row, column=col, value=val)
    c.font = font; c.border = thin_border; c.alignment = Alignment(horizontal='right' if isinstance(val, (int, float)) else 'left')
    if fmt: c.number_format = fmt
    if fill: c.fill = fill
    return c

def section(ws, row, col, val, cols=7):
    for i in range(cols):
        c = ws.cell(row=row, column=col+i)
        c.fill = SECTION_FILL; c.border = thin_border
    c = ws.cell(row=row, column=col, value=val)
    c.font = SECTION_FONT; c.fill = SECTION_FILL
    return c

def set_widths(ws, widths):
    for i, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w


# ══════════════════════════════════════
# Data
# ══════════════════════════════════════
YEARS = ["Year 1", "Year 2", "Year 3", "Year 4", "Year 5"]

# Revenue build
NEW_LOGOS =     [39, 60, 96, 138, 213]
CHURN_RATE =    [0.03, 0.05, 0.05, 0.06, 0.06]
UPGRADE_RATE =  [0.00, 0.15, 0.25, 0.35, 0.35]
PRICE_INC =     [0.00, 0.00, 0.03, 0.03, 0.05]
NEW_ACV =       [77, 95, 112, 130, 145]  # $K
BEG_ARR =       [0, 3000, 10800, 26500, 53000]  # $K
NEW_ARR =       [3000, 5700, 10800, 17900, 30900]
EXPANSION_ARR = [0, 500, 1500, 3400, 6700]
UPGRADE_ARR =   [0, 1400, 3200, 5600, 7400]
PRICE_ARR =     [0, 0, 800, 1600, 2700]
CHURN_ARR =     [90, 300, 600, 1500, 3200]
END_ARR =       [3000, 10800, 26500, 53000, 97500]
REVENUE =       [1500, 6900, 18700, 39800, 75300]  # recognized

# Customers
BEG_CUST =   [0, 39, 96, 186, 315]
CHURN_CUST = [0, 3, 6, 9, 19]
END_CUST =   [39, 96, 186, 315, 565]

# COGS
INFRA =         [96, 180, 420, 780, 1200]
AI_TOKENS =     [24, 72, 168, 336, 600]
SN_LICENSES =   [12, 24, 36, 48, 60]
PARTNER_COMM =  [0, 66, 261, 726, 1650]
TOTAL_COGS =    [132, 356, 940, 2010, 3810]
GROSS_PROFIT =  [1368, 6544, 17760, 37790, 71490]
GROSS_MARGIN =  [0.912, 0.949, 0.950, 0.950, 0.949]

# OpEx
ENG_COST =      [375, 695, 945, 1285, 1740]
CHANNEL_COST =  [0, 140, 280, 425, 570]
SUPPORT_COST =  [0, 0, 130, 260, 385]
FOUNDER_ALLOC = [84, 84, 72, 60, 48]
SN_DEV_ALLOC =  [108, 90, 72, 60, 48]
QA_ALLOC =      [108, 90, 72, 60, 48]
LEGAL =         [90, 15, 10, 25, 10]
MARKETING =     [65, 55, 70, 90, 110]
SELF_SERVE =    [0, 60, 120, 180, 240]
GA_ALLOC =      [150, 690, 1870, 3980, 7530]
TOTAL_OPEX =    [980, 1919, 3639, 6335, 10281]

EBITDA =        [388, 4625, 14121, 31455, 61209]
EBITDA_MARGIN = [0.259, 0.670, 0.755, 0.790, 0.813]
DA =            [100, 150, 200, 250, 300]
EBIT =          [288, 4475, 13921, 31205, 60909]
TAX =           [72, 1119, 3480, 7801, 15227]
NET_INCOME =    [216, 3356, 10441, 23404, 45682]
NET_MARGIN =    [0.144, 0.487, 0.558, 0.588, 0.607]

# Cash flow
DEF_REV =       [750, 1950, 3950, 6650, 11125]
AR_CHANGE =     [125, 450, 983, 1758, 2963]
CAPEX =         [50, 75, 100, 125, 150]
FCF =           [891, 4931, 14008, 28421, 53994]
BEG_CASH =      [0, 891, 5822, 19830, 48251]
END_CASH =      [891, 5822, 19830, 48251, 102245]

# Unit economics
CAC =           [12, 10, 9, 8, 8]  # $K
LTV =           [480, 540, 620, 740, 800]  # $K
LTV_CAC =       [40, 54, 69, 93, 100]
CAC_PAYBACK =   [1.9, 1.3, 1.0, 0.7, 0.6]
NRR =           [0, 1.62, 1.60, 1.55, 1.50]

# Headcount
HEADCOUNT =     [5, 7, 9, 11, 14]
PF_FTE =        [3.6, 4.8, 5.8, 6.4, 7.0]

# Scenarios
SCENARIOS = {
    "Bear":       {"y5_arr": 42000, "y5_cust": 250, "y5_rev": 32000, "y5_ebitda": 21000, "y5_margin": 0.66, "y5_ni": 14000, "y5_hc": 14, "cum_rev": 65000, "cum_ni": 28000},
    "Likely":     {"y5_arr": 97500, "y5_cust": 565, "y5_rev": 75300, "y5_ebitda": 61200, "y5_margin": 0.81, "y5_ni": 45700, "y5_hc": 16, "cum_rev": 142000, "cum_ni": 82000},
    "Bull":       {"y5_arr": 125000, "y5_cust": 730, "y5_rev": 98000, "y5_ebitda": 80000, "y5_margin": 0.82, "y5_ni": 60000, "y5_hc": 18, "cum_rev": 185000, "cum_ni": 110000},
    "Best Case":  {"y5_arr": 155000, "y5_cust": 900, "y5_rev": 121000, "y5_ebitda": 100000, "y5_margin": 0.83, "y5_ni": 75000, "y5_hc": 20, "cum_rev": 230000, "cum_ni": 140000},
}


def build_xlsx():
    wb = Workbook()

    # ── Sheet 1: Cover ──
    ws = wb.active
    ws.title = "Cover"
    ws.cell(row=2, column=2, value="Avennorth Pathfinder").font = Font(name="Calibri", size=24, bold=True, color="1C1917")
    ws.cell(row=3, column=2, value="Five-Year Financial Model").font = Font(name="Calibri", size=16, color="666666")
    ws.cell(row=5, column=2, value="Version 2.0  |  March 2026  |  Confidential").font = DATA_FONT
    ws.cell(row=7, column=2, value="Model Structure:").font = BOLD_FONT
    sheets = ["Assumptions", "Revenue Build", "P&L", "Cash Flow", "Unit Economics", "Headcount", "Scenarios", "Sensitivity", "Portfolio", "Pricing"]
    for i, s in enumerate(sheets):
        ws.cell(row=8+i, column=2, value=f"  {i+1}. {s}").font = DATA_FONT
    ws.cell(row=20, column=2, value="Blue cells = changeable inputs").font = Font(name="Calibri", size=10, color="1565C0")
    ws.cell(row=20, column=2).fill = INPUT_FILL
    set_widths(ws, [4, 50])

    # ── Sheet 2: Assumptions ──
    ws2 = wb.create_sheet("Assumptions")
    ws2.cell(row=1, column=1, value="Key Assumptions").font = TITLE_FONT
    set_widths(ws2, [35, 15, 15, 15, 15, 15])
    r = 3
    for label, vals, fmt in [
        ("New logo acquisition", NEW_LOGOS, MONEY_FMT),
        ("Logo churn rate", CHURN_RATE, PCT_FMT),
        ("Std->Pro upgrade rate", UPGRADE_RATE, PCT_FMT),
        ("Annual price increase", PRICE_INC, PCT_FMT),
        ("New logo ACV ($K)", NEW_ACV, MONEY_FMT),
    ]:
        hdr(ws2, r, 1, label) if r == 3 else None
        cell(ws2, r, 1, label, font=BOLD_FONT)
        for j, v in enumerate(vals):
            if r == 3:
                hdr(ws2, 2, j+2, YEARS[j])
            cell(ws2, r, j+2, v, fmt=fmt, fill=INPUT_FILL)
        r += 1

    r += 1
    cell(ws2, r, 1, "XL Professional blended", font=BOLD_FONT)
    cell(ws2, r, 2, 350000, fmt=MONEY_FMT, fill=INPUT_FILL)
    r += 1
    cell(ws2, r, 1, "G&A allocation (% of revenue)", font=BOLD_FONT)
    cell(ws2, r, 2, 0.10, fmt=PCT_FMT, fill=INPUT_FILL)
    r += 1
    cell(ws2, r, 1, "Tax rate", font=BOLD_FONT)
    cell(ws2, r, 2, 0.25, fmt=PCT_FMT, fill=INPUT_FILL)
    r += 1
    cell(ws2, r, 1, "Existing team allocation ($K/yr)", font=BOLD_FONT)
    cell(ws2, r, 2, 420, fmt=MONEY_FMT, fill=INPUT_FILL)
    for j in range(5):
        hdr(ws2, 2, j+2, YEARS[j])

    # ── Sheet 3: Revenue Build ──
    ws3 = wb.create_sheet("Revenue Build")
    ws3.cell(row=1, column=1, value="Revenue Waterfall ($K)").font = TITLE_FONT
    set_widths(ws3, [30, 15, 15, 15, 15, 15])
    for j in range(5):
        hdr(ws3, 2, j+2, YEARS[j])

    rows3 = [
        ("CUSTOMERS", None, None, True),
        ("Beginning customers", BEG_CUST, MONEY_FMT, False),
        ("+ New logos", NEW_LOGOS, MONEY_FMT, False),
        ("- Churned logos", CHURN_CUST, MONEY_FMT, False),
        ("Ending customers", END_CUST, MONEY_FMT, True),
        ("Logo churn rate", CHURN_RATE, PCT_FMT, False),
        ("", None, None, False),
        ("ARR WATERFALL ($K)", None, None, True),
        ("Beginning ARR", BEG_ARR, MONEY_FMT, False),
        ("+ New logo ARR", NEW_ARR, MONEY_FMT, False),
        ("+ Expansion ARR (tier)", EXPANSION_ARR, MONEY_FMT, False),
        ("+ Std->Pro upgrade ARR", UPGRADE_ARR, MONEY_FMT, False),
        ("+ Price increases", PRICE_ARR, MONEY_FMT, False),
        ("- Churned ARR", CHURN_ARR, MONEY_FMT, False),
        ("Ending ARR", END_ARR, MONEY_FMT, True),
        ("Net new ARR", [END_ARR[i]-BEG_ARR[i] for i in range(5)], MONEY_FMT, True),
        ("Net Revenue Retention", NRR, PCT_FMT, False),
        ("", None, None, False),
        ("RECOGNIZED REVENUE ($K)", None, None, True),
        ("Revenue (midpoint)", REVENUE, MONEY_FMT, True),
        ("Revenue growth %", [0, 3.40, 1.64, 1.09, 0.82], PCT_FMT, False),
    ]
    r = 3
    for label, vals, fmt, is_bold in rows3:
        f = BOLD_FONT if is_bold else DATA_FONT
        fl = TOTAL_FILL if is_bold and vals else (SECTION_FILL if is_bold and not vals else None)
        cell(ws3, r, 1, label, font=f, fill=fl)
        if vals:
            for j, v in enumerate(vals):
                cell(ws3, r, j+2, v, fmt=fmt, font=f, fill=fl)
        r += 1

    # ── Sheet 4: P&L ──
    ws4 = wb.create_sheet("P&L")
    ws4.cell(row=1, column=1, value="Profit & Loss Statement ($K) -- Likely Case").font = TITLE_FONT
    set_widths(ws4, [35, 15, 15, 15, 15, 15])
    for j in range(5):
        hdr(ws4, 2, j+2, YEARS[j])

    pnl_rows = [
        ("Revenue", REVENUE, MONEY_FMT, True, TOTAL_FILL),
        ("", None, None, False, None),
        ("COST OF REVENUE", None, None, True, SECTION_FILL),
        ("Cloud infrastructure", [-x for x in INFRA], MONEY_FMT, False, None),
        ("Claude API tokens", [-x for x in AI_TOKENS], MONEY_FMT, False, None),
        ("ServiceNow licenses", [-x for x in SN_LICENSES], MONEY_FMT, False, None),
        ("Partner commissions", [-x for x in PARTNER_COMM], MONEY_FMT, False, None),
        ("Total COGS", [-x for x in TOTAL_COGS], MONEY_FMT, True, TOTAL_FILL),
        ("Gross Profit", GROSS_PROFIT, MONEY_FMT, True, TOTAL_FILL),
        ("Gross Margin", GROSS_MARGIN, PCT_FMT, True, TOTAL_FILL),
        ("", None, None, False, None),
        ("OPERATING EXPENSES", None, None, True, SECTION_FILL),
        ("Engineering (dedicated)", [-x for x in ENG_COST], MONEY_FMT, False, None),
        ("Channel / GTM", [-x for x in CHANNEL_COST], MONEY_FMT, False, None),
        ("Support / CS", [-x for x in SUPPORT_COST], MONEY_FMT, False, None),
        ("Founder allocation", [-x for x in FOUNDER_ALLOC], MONEY_FMT, False, None),
        ("SN Developer allocation", [-x for x in SN_DEV_ALLOC], MONEY_FMT, False, None),
        ("QA/DevOps allocation", [-x for x in QA_ALLOC], MONEY_FMT, False, None),
        ("Legal + patents", [-x for x in LEGAL], MONEY_FMT, False, None),
        ("Marketing + partner", [-x for x in MARKETING], MONEY_FMT, False, None),
        ("Self-serve onboarding", [-x for x in SELF_SERVE], MONEY_FMT, False, None),
        ("G&A allocation (10%)", [-x for x in GA_ALLOC], MONEY_FMT, False, None),
        ("Total OpEx", [-x for x in TOTAL_OPEX], MONEY_FMT, True, TOTAL_FILL),
        ("", None, None, False, None),
        ("EBITDA", EBITDA, MONEY_FMT, True, TOTAL_FILL),
        ("EBITDA Margin", EBITDA_MARGIN, PCT_FMT, True, TOTAL_FILL),
        ("D&A", [-x for x in DA], MONEY_FMT, False, None),
        ("EBIT", EBIT, MONEY_FMT, True, None),
        ("Tax (25%)", [-x for x in TAX], MONEY_FMT, False, None),
        ("Net Income", NET_INCOME, MONEY_FMT, True, TOTAL_FILL),
        ("Net Margin", NET_MARGIN, PCT_FMT, True, TOTAL_FILL),
    ]
    r = 3
    for label, vals, fmt, is_bold, fl in pnl_rows:
        f = TOTAL_FONT if is_bold and fl == TOTAL_FILL else (SECTION_FONT if fl == SECTION_FILL else (BOLD_FONT if is_bold else DATA_FONT))
        cell(ws4, r, 1, label, font=f, fill=fl)
        if vals:
            for j, v in enumerate(vals):
                cell(ws4, r, j+2, v, fmt=fmt, font=f, fill=fl)
        elif fl:
            for j in range(5):
                ws4.cell(row=r, column=j+2).fill = fl
                ws4.cell(row=r, column=j+2).border = thin_border
        r += 1

    # ── Sheet 5: Cash Flow ──
    ws5 = wb.create_sheet("Cash Flow")
    ws5.cell(row=1, column=1, value="Cash Flow Statement ($K) -- Likely Case").font = TITLE_FONT
    set_widths(ws5, [35, 15, 15, 15, 15, 15])
    for j in range(5):
        hdr(ws5, 2, j+2, YEARS[j])

    cf_rows = [
        ("Net Income", NET_INCOME, MONEY_FMT, False, None),
        ("+ D&A", DA, MONEY_FMT, False, None),
        ("+ Change in deferred revenue", DEF_REV, MONEY_FMT, False, None),
        ("- Change in A/R", [-x for x in AR_CHANGE], MONEY_FMT, False, None),
        ("- Capex", [-x for x in CAPEX], MONEY_FMT, False, None),
        ("Free Cash Flow", FCF, MONEY_FMT, True, TOTAL_FILL),
        ("", None, None, False, None),
        ("Beginning cash", BEG_CASH, MONEY_FMT, False, None),
        ("+ Free cash flow", FCF, MONEY_FMT, False, None),
        ("Ending cash", END_CASH, MONEY_FMT, True, TOTAL_FILL),
    ]
    r = 3
    for label, vals, fmt, is_bold, fl in cf_rows:
        f = TOTAL_FONT if fl == TOTAL_FILL else (BOLD_FONT if is_bold else DATA_FONT)
        cell(ws5, r, 1, label, font=f, fill=fl)
        if vals:
            for j, v in enumerate(vals):
                cell(ws5, r, j+2, v, fmt=fmt, font=f, fill=fl)
        r += 1

    # ── Sheet 6: Unit Economics ──
    ws6 = wb.create_sheet("Unit Economics")
    ws6.cell(row=1, column=1, value="Unit Economics -- Likely Case").font = TITLE_FONT
    set_widths(ws6, [30, 15, 15, 15, 15, 15])
    for j in range(5):
        hdr(ws6, 2, j+2, YEARS[j])

    ue_rows = [
        ("Blended ACV ($K)", NEW_ACV, MONEY_FMT),
        ("Gross margin", GROSS_MARGIN, PCT_FMT),
        ("CAC ($K)", CAC, MONEY_FMT),
        ("LTV ($K, 5-yr GM)", LTV, MONEY_FMT),
        ("LTV / CAC", LTV_CAC, RATIO_FMT),
        ("CAC payback (months)", CAC_PAYBACK, MONTHS_FMT),
        ("NRR", NRR, PCT_FMT),
        ("ARR per employee ($K)", [int(END_ARR[i]/HEADCOUNT[i]) for i in range(5)], MONEY_FMT),
        ("Revenue per FTE ($K)", [int(REVENUE[i]/PF_FTE[i]) for i in range(5)], MONEY_FMT),
    ]
    r = 3
    for label, vals, fmt in ue_rows:
        cell(ws6, r, 1, label, font=BOLD_FONT)
        for j, v in enumerate(vals):
            cell(ws6, r, j+2, v, fmt=fmt)
        r += 1

    # ── Sheet 7: Headcount ──
    ws7 = wb.create_sheet("Headcount")
    ws7.cell(row=1, column=1, value="Headcount Plan").font = TITLE_FONT
    set_widths(ws7, [30, 12, 12, 12, 12, 12, 15, 20])
    cols7 = ["Role"] + YEARS + ["Loaded $K", "Type"]
    for j, c in enumerate(cols7):
        hdr(ws7, 2, j+1, c)

    roles = [
        ("Founder / Architect", [1,1,1,1,1], 0, "Shared 40%->20%"),
        ("SN Developer", [1,1,1,1,1], 0, "Shared 60%->25%"),
        ("QA / DevOps", [1,1,1,1,1], 0, "Shared 60%->25%"),
        ("Go / Systems Engineer", [1,1,1,1,1], 190, "PF dedicated"),
        ("Python / AI Engineer", [1,1,1,1,1], 185, "PF dedicated"),
        ("Go / Platform Engineer", [0,1,1,1,1], 180, "PF dedicated"),
        ("Channel Manager", [0,1,1,1,1], 140, "Shared 100%->50%"),
        ("Support Engineer", [0,0,1,1,1], 130, "PF dedicated"),
        ("Content / Dev Marketing", [0,0,1,1,1], 120, "Shared"),
        ("Senior AI/ML Engineer", [0,0,0,1,1], 210, "Shared"),
        ("2nd Support Engineer", [0,0,0,1,1], 130, "PF dedicated"),
        ("2nd Channel Manager", [0,0,0,0,1], 145, "Shared"),
        ("Platform Engineer", [0,0,0,0,1], 185, "PF dedicated"),
        ("CSM", [0,0,0,0,1], 125, "PF dedicated"),
    ]
    r = 3
    for name, counts, cost, typ in roles:
        cell(ws7, r, 1, name, font=DATA_FONT)
        for j, c in enumerate(counts):
            cell(ws7, r, j+2, c if c else "", fmt=MONEY_FMT)
        cell(ws7, r, 7, cost if cost else "Existing", fmt=MONEY_FMT)
        cell(ws7, r, 8, typ)
        r += 1
    r += 1
    cell(ws7, r, 1, "Total Headcount", font=TOTAL_FONT, fill=TOTAL_FILL)
    for j, h in enumerate(HEADCOUNT):
        cell(ws7, r, j+2, h, fmt=MONEY_FMT, font=TOTAL_FONT, fill=TOTAL_FILL)
    r += 1
    cell(ws7, r, 1, "PF FTE-equivalent", font=BOLD_FONT)
    for j, f in enumerate(PF_FTE):
        cell(ws7, r, j+2, f, fmt='0.0')

    # ── Sheet 8: Scenarios ──
    ws8 = wb.create_sheet("Scenarios")
    ws8.cell(row=1, column=1, value="Four-Scenario Analysis (Year 5)").font = TITLE_FONT
    set_widths(ws8, [30, 18, 18, 18, 18])
    scen_names = list(SCENARIOS.keys())
    for j, s in enumerate(scen_names):
        hdr(ws8, 2, j+2, s)

    scen_metrics = [
        ("Y5 Ending ARR ($K)", "y5_arr", MONEY_FMT),
        ("Y5 Revenue ($K)", "y5_rev", MONEY_FMT),
        ("Y5 Customers", "y5_cust", MONEY_FMT),
        ("Y5 EBITDA ($K)", "y5_ebitda", MONEY_FMT),
        ("Y5 EBITDA Margin", "y5_margin", PCT_FMT),
        ("Y5 Net Income ($K)", "y5_ni", MONEY_FMT),
        ("Y5 Headcount", "y5_hc", MONEY_FMT),
        ("Y5 ARR/Employee ($K)", None, MONEY_FMT),
        ("5-Year Cumulative Revenue ($K)", "cum_rev", MONEY_FMT),
        ("5-Year Cumulative Net Income ($K)", "cum_ni", MONEY_FMT),
    ]
    r = 3
    for label, key, fmt in scen_metrics:
        cell(ws8, r, 1, label, font=BOLD_FONT)
        for j, s in enumerate(scen_names):
            if key:
                v = SCENARIOS[s][key]
            else:
                v = int(SCENARIOS[s]["y5_arr"] / SCENARIOS[s]["y5_hc"])
            is_likely = s == "Likely"
            cell(ws8, r, j+2, v, fmt=fmt, fill=TOTAL_FILL if is_likely else None, font=TOTAL_FONT if is_likely else DATA_FONT)
        r += 1

    # ── Sheet 9: Sensitivity ──
    ws9 = wb.create_sheet("Sensitivity")
    ws9.cell(row=1, column=1, value="Sensitivity Analysis -- Y5 ARR Impact ($K)").font = TITLE_FONT
    set_widths(ws9, [30, 18, 18, 18])
    for j, h in enumerate(["-20% Change", "Base", "+20% Change"]):
        hdr(ws9, 2, j+2, h)

    sens = [
        ("New logo acquisition rate", 62000, 84000, 106000),
        ("Std->Pro upgrade rate", 72000, 84000, 96000),
        ("Logo churn rate", 91000, 84000, 77000),
        ("Blended ACV", 67000, 84000, 101000),
        ("Price increase rate", 79000, 84000, 89000),
    ]
    for i, (label, lo, base, hi) in enumerate(sens, 3):
        cell(ws9, i, 1, label, font=BOLD_FONT)
        cell(ws9, i, 2, lo, fmt=MONEY_FMT)
        cell(ws9, i, 3, base, fmt=MONEY_FMT, fill=TOTAL_FILL, font=TOTAL_FONT)
        cell(ws9, i, 4, hi, fmt=MONEY_FMT)

    # ── Sheet 10: Portfolio ──
    ws10 = wb.create_sheet("Portfolio")
    ws10.cell(row=1, column=1, value="Avennorth Portfolio View (Illustrative)").font = TITLE_FONT
    set_widths(ws10, [25, 15, 15, 15, 15, 15, 25])
    for j in range(5):
        hdr(ws10, 2, j+2, YEARS[j])
    hdr(ws10, 2, 7, "Notes")

    portfolio = [
        ("PRODUCT ARR ($K)", None, True),
        ("Pathfinder", END_ARR, False),
        ("Bearing", [0, 500, 3000, 8000, 15000], False),
        ("Contour", [0, 0, 0, 2000, 8000], False),
        ("Vantage", [0, 0, 0, 0, 3000], False),
        ("CoreX Consulting", [5000, 5500, 6000, 6500, 7000], False),
        ("Compass Platform", [1500, 2000, 2500, 3000, 4000], False),
        ("Total Avennorth ARR", [9500, 18200, 36000, 67500, 121000], True),
        ("", None, False),
        ("SHARED COST ALLOCATION", None, True),
        ("Shared AI Engine -> PF", [0.80, 0.70, 0.60, 0.50, 0.45], False),
        ("Compass channel -> PF", [1.00, 0.80, 0.65, 0.55, 0.50], False),
        ("Founder -> PF", [0.40, 0.40, 0.30, 0.25, 0.20], False),
        ("SN Dev -> PF", [0.60, 0.50, 0.40, 0.30, 0.25], False),
        ("QA/DevOps -> PF", [0.60, 0.50, 0.40, 0.30, 0.25], False),
        ("G&A -> PF", [1.00, 0.60, 0.45, 0.35, 0.30], False),
    ]
    notes = {
        "Pathfinder": "Modeled (this workbook)",
        "Bearing": "Assessment-driven",
        "Contour": "Launches Y3",
        "Vantage": "Launches Y4",
        "CoreX Consulting": "Existing, grows with portfolio",
        "Compass Platform": "Existing SaaS",
        "Total Avennorth ARR": "Pathfinder ~70% of Y5",
    }
    r = 3
    for label, vals, is_bold in portfolio:
        f = TOTAL_FONT if is_bold and vals else (SECTION_FONT if is_bold and not vals else DATA_FONT)
        fl = TOTAL_FILL if is_bold and vals else (SECTION_FILL if is_bold and not vals else None)
        cell(ws10, r, 1, label, font=f, fill=fl)
        if vals:
            fmt = PCT_FMT if isinstance(vals[0], float) and vals[0] <= 1 else MONEY_FMT
            for j, v in enumerate(vals):
                cell(ws10, r, j+2, v, fmt=fmt, font=f, fill=fl)
        elif fl:
            for j in range(5):
                ws10.cell(row=r, column=j+2).fill = fl
                ws10.cell(row=r, column=j+2).border = thin_border
        if label in notes:
            cell(ws10, r, 7, notes[label])
        r += 1

    # ── Sheet 11: Pricing ──
    ws11 = wb.create_sheet("Pricing")
    ws11.cell(row=1, column=1, value="Pricing Structure").font = TITLE_FONT
    set_widths(ws11, [20, 20, 25, 25])

    for j, h in enumerate(["Tier", "Managed Nodes", "Standard (starting at)", "Professional (starting at)"]):
        hdr(ws11, 3, j+1, h)
    pricing = [
        ("S", "Up to 500", 50000, 100000),
        ("M", "501 - 2,000", 90000, 175000),
        ("L", "2,001 - 5,000", 150000, 250000),
        ("XL", "5,001+", 200000, "Custom"),
    ]
    for i, (tier, nodes, std, pro) in enumerate(pricing, 4):
        cell(ws11, i, 1, tier, font=BOLD_FONT)
        cell(ws11, i, 2, nodes)
        cell(ws11, i, 3, std, fmt='$#,##0')
        cell(ws11, i, 4, pro, fmt='$#,##0' if isinstance(pro, int) else None)

    r = 10
    cell(ws11, r, 1, "PACKAGES", font=SECTION_FONT, fill=SECTION_FILL)
    for j in range(3): ws11.cell(row=r, column=j+2).fill = SECTION_FILL
    r += 1
    cell(ws11, r, 1, "Standard", font=BOLD_FONT)
    cell(ws11, r, 2, "Pathfinder Discovery Engine + CMDB Ops")
    r += 1
    cell(ws11, r, 1, "Professional", font=BOLD_FONT)
    cell(ws11, r, 2, "Standard + Integration Intelligence + Service Map Intelligence")
    r += 2
    cell(ws11, r, 1, "MANAGED NODE", font=SECTION_FONT, fill=SECTION_FILL)
    for j in range(3): ws11.cell(row=r, column=j+2).fill = SECTION_FILL
    r += 1
    cell(ws11, r, 1, "Counted", font=BOLD_FONT)
    cell(ws11, r, 2, "Servers, VMs, cloud instances running Pathfinder agent")
    r += 1
    cell(ws11, r, 1, "Free (not counted)", font=BOLD_FONT)
    cell(ws11, r, 2, "Network devices, desktops, agentless discovery")

    # Save
    xlsx_path = os.path.join(ICLOUD_FINANCIAL, "Pathfinder_Financial_Model.xlsx")
    wb.save(xlsx_path)
    size_kb = os.path.getsize(xlsx_path) // 1024
    print(f"Generated XLSX: {xlsx_path} ({size_kb} KB)")


def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    LIME = RGBColor(0xC8, 0xFF, 0x00)
    DARK = RGBColor(0x0C, 0x0C, 0x18)
    WHITE = RGBColor(0xF0, 0xF0, 0xF8)
    GRAY = RGBColor(0x6A, 0x6A, 0x84)
    TEAL = RGBColor(0x00, 0xD4, 0xAA)
    ORANGE = RGBColor(0xFF, 0x8C, 0x40)

    def dark_slide(title, subtitle=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = DARK
        tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(1))
        tf = tx.text_frame
        p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(32); p.font.bold = True; p.font.color.rgb = WHITE
        if subtitle:
            p2 = tf.add_paragraph(); p2.text = subtitle; p2.font.size = Pt(16); p2.font.color.rgb = GRAY
        return slide

    def txt(slide, l, t, w, h, text, sz=14, color=WHITE, bold=False):
        tx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tx.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold
        return tf

    def kpi(slide, l, t, value, label, color=LIME):
        txt(slide, l, t, 2.5, 0.5, value, sz=28, color=color, bold=True)
        txt(slide, l, t+0.5, 2.5, 0.3, label, sz=11, color=GRAY)

    # Slide 1: Title
    s = dark_slide("Avennorth Pathfinder", "Five-Year Business Case  |  v2.0  |  March 2026")
    txt(s, 0.8, 2.2, 10, 0.5, "CMDB-first integration discovery + AI intelligence for ServiceNow", sz=18, color=GRAY)
    txt(s, 0.8, 3.2, 10, 0.4, "Two packages: Standard ($50K+) and Professional ($100K+)", sz=16, color=TEAL)
    txt(s, 0.8, 3.8, 10, 0.4, "Penetration pricing is the base plan. Profitable in Year 1.", sz=16, color=LIME)

    # Slide 2: Key Metrics
    s = dark_slide("Investor Summary -- Likely Case (Year 5)")
    kpi(s, 0.8, 1.8, "$84M", "Y5 ARR")
    kpi(s, 3.3, 1.8, "$66M", "Y5 Revenue", TEAL)
    kpi(s, 5.8, 1.8, "80%", "EBITDA Margin", TEAL)
    kpi(s, 8.3, 1.8, "$39M", "Net Income", LIME)
    kpi(s, 10.8, 1.8, "$89M", "Cash Position", LIME)

    kpi(s, 0.8, 3.5, "546", "Customers", TEAL)
    kpi(s, 3.3, 3.5, "145%", "NRR", TEAL)
    kpi(s, 5.8, 3.5, "100x", "LTV/CAC", LIME)
    kpi(s, 8.3, 3.5, "~16", "Headcount", GRAY)
    kpi(s, 10.8, 3.5, "$5.3M", "ARR/Employee", LIME)

    txt(s, 0.8, 5.5, 11, 1,
        "5-Year Cumulative: $128M revenue  |  $73M net income  |  $89M ending cash  |  $0 external funding required",
        sz=16, color=WHITE)

    # Slide 3: Revenue Waterfall
    s = dark_slide("Revenue Build", "ARR waterfall -- new logos + expansion + upgrades - churn")
    txt(s, 0.8, 1.8, 11, 3,
        "             Y1          Y2          Y3          Y4          Y5\n"
        "Begin      $0          $3.0M     $10.2M    $24.5M    $48.0M\n"
        "+ New      $3.0M      $5.7M     $10.8M    $17.9M    $30.9M\n"
        "+ Expand   $0          $0.5M      $1.5M      $3.4M      $6.7M\n"
        "+ Upgrade  $0          $1.4M      $3.2M      $5.6M      $7.4M\n"
        "+ Price      $0          $0            $0.7M      $1.4M      $2.4M\n"
        "- Churn    ($0.2M)   ($0.4M)   ($1.9M)   ($4.8M)  ($11.4M)\n"
        "End ARR   $3.0M    $10.2M    $24.5M    $48.0M    $84.0M",
        sz=16, color=WHITE)
    txt(s, 0.8, 5.8, 11, 0.5, "NRR: 157% -> 153% -> 148% -> 145%.  New logo acquisition is the #1 lever.", sz=14, color=TEAL)

    # Slide 4: P&L Summary
    s = dark_slide("P&L Summary", "Gross margin ~95%. EBITDA margin scales from 26% to 80%.")
    txt(s, 0.8, 1.8, 11, 3,
        "                   Y1          Y2          Y3          Y4          Y5\n"
        "Revenue      $1.5M      $6.6M     $17.4M    $36.3M    $66.0M\n"
        "COGS          ($132K)   ($342K)   ($885K)   ($1.9M)   ($3.5M)\n"
        "Gross Profit  $1.4M      $6.3M     $16.5M    $34.4M    $62.5M\n"
        "GM %            91%         95%         95%         95%         95%\n\n"
        "OpEx            ($980K)   ($1.9M)   ($3.5M)   ($6.1M)   ($9.8M)\n"
        "EBITDA        $388K      $4.4M     $13.0M    $28.3M    $52.7M\n"
        "EBITDA %      26%          66%         75%         78%         80%\n\n"
        "Net Income   $216K      $3.2M      $9.6M     $21.0M    $39.3M",
        sz=15, color=WHITE)

    # Slide 5: Scenarios
    s = dark_slide("Scenario Analysis", "Bear / Likely / Bull / Best Case -- Y5 metrics")
    labels = ["Y5 ARR", "Y5 Revenue", "Customers", "EBITDA Margin", "Net Income", "Headcount"]
    vals = [
        ["$33M", "$84M", "$110M", "$135M"],
        ["$25M", "$66M", "$86M", "$106M"],
        ["217", "546", "710", "875"],
        ["68%", "80%", "81%", "83%"],
        ["$11M", "$39M", "$52M", "$65M"],
        ["14", "16", "18", "20"],
    ]
    names = ["Bear", "Likely (Base)", "Bull", "Best Case"]
    for j, name in enumerate(names):
        x = 3 + j * 2.6
        txt(s, x, 1.5, 2.4, 0.4, name, sz=16,
            color=LIME if name == "Likely (Base)" else (ORANGE if name == "Bear" else WHITE), bold=True)
    for i, label in enumerate(labels):
        txt(s, 0.8, 2.2 + i * 0.7, 2, 0.4, label, sz=13, color=GRAY, bold=True)
        for j in range(4):
            txt(s, 3 + j * 2.6, 2.2 + i * 0.7, 2.4, 0.4, vals[i][j], sz=13,
                color=LIME if j == 1 else WHITE)

    # Slide 6: Portfolio
    s = dark_slide("Avennorth Portfolio", "Pathfinder is ~70% of Y5 ARR. Shared costs decline as portfolio grows.")
    txt(s, 0.8, 1.8, 5.5, 0.4, "PORTFOLIO ARR (Y5)", sz=18, color=LIME, bold=True)
    txt(s, 0.8, 2.4, 5.5, 2.5,
        "Pathfinder        $84M   (modeled)\n"
        "Bearing             $15M   (assessment-driven)\n"
        "Contour              $8M    (launches Y3)\n"
        "Vantage              $3M    (launches Y4)\n"
        "CoreX                  $7M    (existing consulting)\n"
        "Compass              $4M    (platform SaaS)\n"
        "Total                  $121M",
        sz=15, color=WHITE)

    txt(s, 7, 1.8, 5, 0.4, "SHARED COST LEVERAGE", sz=18, color=TEAL, bold=True)
    txt(s, 7, 2.4, 5, 2.5,
        "As portfolio grows, Pathfinder's share\n"
        "of fixed costs drops ~50% (Y1 to Y5):\n\n"
        "Shared AI Engine:  80% -> 35%\n"
        "Compass channel:  100% -> 40%\n"
        "Founder time:         40% -> 20%\n"
        "SN Dev / QA:          60% -> 25%\n"
        "G&A overhead:        100% -> 30%",
        sz=14, color=WHITE)

    txt(s, 0.8, 5.8, 11, 0.5,
        "Portfolio flywheel: Pathfinder -> Bearing -> Contour -> Vantage.  Each product cross-sells the next.",
        sz=14, color=TEAL)

    # Slide 7: Unit Economics
    s = dark_slide("Unit Economics", "LTV/CAC scales from 40x to 100x. CAC payback < 2 months.")
    kpi(s, 0.8, 1.8, "$8K", "Y5 CAC", TEAL)
    kpi(s, 3.3, 1.8, "$800K", "Y5 LTV (5yr GM)", LIME)
    kpi(s, 5.8, 1.8, "100x", "LTV/CAC", LIME)
    kpi(s, 8.3, 1.8, "0.6 mo", "CAC Payback", TEAL)
    kpi(s, 10.8, 1.8, "145%", "Y5 NRR", LIME)

    txt(s, 0.8, 3.5, 11, 2,
        "No direct sales team. CAC is dominated by channel enablement and self-serve tooling.\n\n"
        "CAC composition: channel manager ($2.7K) + marketing ($0.5K) + self-serve ($1.1K) + other ($3.7K)\n\n"
        "NRR driven by: Std->Pro upgrades (35% rate) + node tier expansion (20%) + price increases (5%)\n\n"
        "Cohort economics: Y1 cohort (39 logos) generates $11.2M over 3 years at 95% gross margin.",
        sz=14, color=WHITE)

    # Slide 8: Cash & Funding
    s = dark_slide("Cash Flow & Funding", "Bootstrappable. Profitable in Year 1. $89M ending cash at Y5.")
    txt(s, 0.8, 1.8, 5.5, 3,
        "CASH FLOW ($K)\n\n"
        "         Y1       Y2       Y3        Y4        Y5\n"
        "FCF    $891   $4.6M  $12.4M  $25.5M  $46.0M\n"
        "Cash   $891   $5.5M  $17.9M  $43.4M  $89.4M",
        sz=16, color=WHITE)

    txt(s, 7, 1.8, 5, 0.4, "FUNDING OPTIONS", sz=18, color=TEAL, bold=True)
    txt(s, 7, 2.4, 5, 2.5,
        "Bootstrap (recommended):\n"
        "  $0 external. $610K incremental Y1.\n"
        "  Absorbed by consulting revenue.\n\n"
        "Seed ($1.5-2M, optional):\n"
        "  +2 engineers, channel manager,\n"
        "  conference presence. Accelerates\n"
        "  Bull/Best case probability.\n\n"
        "Series A (only if Bull/Best):\n"
        "  $8-12M expansion capital.",
        sz=14, color=WHITE)

    pptx_path = os.path.join(ICLOUD_PRESENTATIONS, "Pathfinder_Business_Case.pptx")
    prs.save(pptx_path)
    size_kb = os.path.getsize(pptx_path) // 1024
    print(f"Generated PPTX: {pptx_path} ({size_kb} KB)")


if __name__ == "__main__":
    print("=" * 60)
    print("Avennorth Pathfinder -- Financial Model Generator")
    print("=" * 60)
    build_xlsx()
    build_pptx()
    print("Done.")
