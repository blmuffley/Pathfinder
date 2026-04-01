#!/usr/bin/env python3
"""Generate all iCloud output files for Pathfinder.

Produces:
- DOCX files in iCloud/Docs/ (from updated MD files)
- XLSX financial model in iCloud/Financial/
- PPTX summary in iCloud/Presentations/
- PDF files in iCloud/Docs/

Run: python3 docs/generate_icloud_outputs.py
"""

import os
import re
from datetime import datetime

# ── Paths ──
DOCS_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.dirname(DOCS_DIR)
ICLOUD_BASE = os.path.expanduser(
    "~/Library/Mobile Documents/com~apple~CloudDocs/Projects/Avennorth/Solutions/Pathfinder"
)
ICLOUD_DOCS = os.path.join(ICLOUD_BASE, "Docs")
ICLOUD_FINANCIAL = os.path.join(ICLOUD_BASE, "Financial")
ICLOUD_PRESENTATIONS = os.path.join(ICLOUD_BASE, "Presentations")

for d in [ICLOUD_DOCS, ICLOUD_FINANCIAL, ICLOUD_PRESENTATIONS]:
    os.makedirs(d, exist_ok=True)
    os.makedirs(os.path.join(d, "markdown"), exist_ok=True) if d == ICLOUD_DOCS else None

# ── MD files to convert ──
MD_FILES = [
    ("guides/five-year-business-case.md", "Pathfinder Five-Year Business Case"),
    ("guides/partner-enablement-guide.md", "Pathfinder Partner Enablement Guide"),
    ("customer-facing/solution-brief.md", "Pathfinder Solution Brief"),
    ("customer-facing/faq.md", "Pathfinder FAQ"),
    ("internal/product-strategy.md", "Pathfinder Product Strategy"),
    ("architecture/04-portfolio-architecture.md", "Pathfinder Portfolio Architecture"),
    ("reference/capacity-planning.md", "Pathfinder Capacity Planning Guide"),
    ("guides/installation-guide.md", "Pathfinder Installation Guide"),
    ("guides/operations-runbook.md", "Pathfinder Operations Runbook"),
    ("guides/implementation-playbook.md", "Pathfinder Implementation Playbook"),
    ("guides/user-guide-workspace.md", "Pathfinder Workspace User Guide"),
    ("reference/api-reference.md", "Pathfinder API Reference"),
    ("reference/security-architecture.md", "Pathfinder Security Architecture"),
    ("reference/data-dictionary.md", "Pathfinder Data Dictionary"),
    ("prototypes/DEMO-SCRIPT.md", "Pathfinder Demo Script"),
]


# ════════════════════════════════════════════
# DOCX Generation
# ════════════════════════════════════════════

def generate_docx_files():
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT

    GRAY = RGBColor(0x78, 0x71, 0x6C)

    def set_style(doc):
        style = doc.styles['Normal']
        font = style.font
        font.name = 'Calibri'
        font.size = Pt(10)
        font.color.rgb = RGBColor(0x2D, 0x2D, 0x2D)
        style.paragraph_format.space_after = Pt(4)
        style.paragraph_format.line_spacing = 1.15

    def add_header_footer(doc, title):
        section = doc.sections[0]
        header = section.header
        hp = header.paragraphs[0]
        hp.text = "AVENNORTH  |  CONFIDENTIAL"
        hp.style.font.size = Pt(8)
        hp.style.font.color.rgb = GRAY
        hp.alignment = WD_ALIGN_PARAGRAPH.LEFT
        footer = section.footer
        fp = footer.paragraphs[0]
        fp.text = f"Avennorth Pathfinder  |  {title}"
        fp.style.font.size = Pt(8)
        fp.style.font.color.rgb = GRAY
        fp.alignment = WD_ALIGN_PARAGRAPH.CENTER

    def parse_md_to_docx(md_path, docx_path, title):
        with open(md_path, 'r', encoding='utf-8') as f:
            lines = f.readlines()

        doc = Document()
        set_style(doc)
        add_header_footer(doc, title)

        in_code_block = False
        in_table = False
        table_rows = []

        def flush_table():
            nonlocal in_table, table_rows
            if not table_rows:
                return
            headers = [c.strip() for c in table_rows[0].split('|')[1:-1]]
            data = []
            for row in table_rows[2:]:
                cells = [c.strip() for c in row.split('|')[1:-1]]
                if cells:
                    data.append(cells)
            if headers:
                ncols = len(headers)
                tbl = doc.add_table(rows=1 + len(data), cols=ncols)
                tbl.style = 'Table Grid'
                tbl.alignment = WD_TABLE_ALIGNMENT.LEFT
                for i, h in enumerate(headers):
                    cell = tbl.rows[0].cells[i]
                    cell.text = h
                    for p in cell.paragraphs:
                        for r in p.runs:
                            r.font.bold = True
                            r.font.size = Pt(8)
                for ri, row_data in enumerate(data):
                    for ci in range(min(ncols, len(row_data))):
                        cell = tbl.rows[ri + 1].cells[ci]
                        cell.text = row_data[ci]
                        for p in cell.paragraphs:
                            for r in p.runs:
                                r.font.size = Pt(8)
            in_table = False
            table_rows = []

        for line in lines:
            stripped = line.rstrip('\n')
            if stripped.startswith('```'):
                if in_table:
                    flush_table()
                in_code_block = not in_code_block
                continue
            if in_code_block:
                p = doc.add_paragraph()
                run = p.add_run(stripped)
                run.font.name = 'Consolas'
                run.font.size = Pt(8)
                run.font.color.rgb = RGBColor(0x40, 0x40, 0x40)
                p.paragraph_format.space_before = Pt(0)
                p.paragraph_format.space_after = Pt(0)
                continue
            if '|' in stripped and stripped.strip().startswith('|'):
                if not in_table:
                    in_table = True
                    table_rows = []
                table_rows.append(stripped)
                continue
            elif in_table:
                flush_table()
            if stripped.startswith('# '):
                doc.add_heading(stripped[2:].strip(), level=1)
                continue
            if stripped.startswith('## '):
                doc.add_heading(stripped[3:].strip(), level=2)
                continue
            if stripped.startswith('### '):
                doc.add_heading(stripped[4:].strip(), level=3)
                continue
            if stripped.startswith('#### '):
                doc.add_heading(stripped[5:].strip(), level=4)
                continue
            if stripped.startswith('---'):
                doc.add_paragraph('_' * 60)
                continue
            if stripped.startswith('- ') or stripped.startswith('* '):
                p = doc.add_paragraph(style='List Bullet')
                text = stripped[2:].strip()
                parts = re.split(r'\*\*(.+?)\*\*', text)
                for i, part in enumerate(parts):
                    run = p.add_run(part)
                    run.font.size = Pt(9)
                    if i % 2 == 1:
                        run.bold = True
                continue
            if stripped.startswith('> '):
                p = doc.add_paragraph()
                p.paragraph_format.left_indent = Inches(0.5)
                run = p.add_run(stripped[2:].strip())
                run.font.size = Pt(9)
                run.italic = True
                continue
            if re.match(r'^\d+\. ', stripped):
                p = doc.add_paragraph(style='List Number')
                text = re.sub(r'^\d+\.\s*', '', stripped)
                p.add_run(text).font.size = Pt(9)
                continue
            if not stripped.strip():
                doc.add_paragraph('')
                continue
            p = doc.add_paragraph()
            text = stripped.strip()
            parts = re.split(r'(\*\*.*?\*\*|`.*?`)', text)
            for part in parts:
                if part.startswith('**') and part.endswith('**'):
                    run = p.add_run(part[2:-2])
                    run.bold = True
                    run.font.size = Pt(10)
                elif part.startswith('`') and part.endswith('`'):
                    run = p.add_run(part[1:-1])
                    run.font.name = 'Consolas'
                    run.font.size = Pt(9)
                    run.font.color.rgb = RGBColor(0x60, 0x60, 0x80)
                else:
                    run = p.add_run(part)
                    run.font.size = Pt(10)

        if in_table:
            flush_table()

        doc.save(docx_path)

    print("Generating DOCX files to iCloud/Docs/...\n")
    md_dir = os.path.join(ICLOUD_DOCS, "markdown")
    os.makedirs(md_dir, exist_ok=True)

    for md_rel, title in MD_FILES:
        md_path = os.path.join(DOCS_DIR, md_rel)
        if not os.path.exists(md_path):
            print(f"  SKIP: {md_rel} (not found)")
            continue
        basename = os.path.splitext(os.path.basename(md_rel))[0]
        docx_path = os.path.join(ICLOUD_DOCS, f"{basename}.docx")
        parse_md_to_docx(md_path, docx_path, title)
        # Copy MD to markdown subfolder
        import shutil
        shutil.copy2(md_path, os.path.join(md_dir, os.path.basename(md_rel)))
        size_kb = os.path.getsize(docx_path) // 1024
        print(f"  OK: {basename}.docx ({size_kb} KB)")

    print()


# ════════════════════════════════════════════
# XLSX Financial Model
# ════════════════════════════════════════════

def generate_xlsx():
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side, numbers
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    LIME = "C8FF00"
    DARK = "0C0C18"
    HEADER_FILL = PatternFill(start_color="1E1E32", end_color="1E1E32", fill_type="solid")
    INPUT_FILL = PatternFill(start_color="1A2A1A", end_color="1A2A1A", fill_type="solid")
    HEADER_FONT = Font(name="Calibri", size=10, bold=True, color="FFFFFF")
    DATA_FONT = Font(name="Calibri", size=10, color="333333")
    TITLE_FONT = Font(name="Calibri", size=14, bold=True, color="1C1917")
    SUBTITLE_FONT = Font(name="Calibri", size=11, bold=True, color="39FF14")
    MONEY_FMT = '#,##0'
    PCT_FMT = '0.0%'
    thin_border = Border(
        left=Side(style='thin'), right=Side(style='thin'),
        top=Side(style='thin'), bottom=Side(style='thin')
    )

    def style_header_row(ws, row, cols):
        for c in range(1, cols + 1):
            cell = ws.cell(row=row, column=c)
            cell.font = HEADER_FONT
            cell.fill = HEADER_FILL
            cell.alignment = Alignment(horizontal='center')
            cell.border = thin_border

    def style_data_cell(ws, row, col, fmt=None):
        cell = ws.cell(row=row, column=col)
        cell.font = DATA_FONT
        cell.border = thin_border
        cell.alignment = Alignment(horizontal='center')
        if fmt:
            cell.number_format = fmt
        return cell

    # ── Sheet 1: Assumptions ──
    ws = wb.active
    ws.title = "Assumptions"
    ws.cell(row=1, column=1, value="Pathfinder Financial Model").font = TITLE_FONT
    ws.cell(row=2, column=1, value=f"Generated {datetime.now().strftime('%Y-%m-%d')}").font = DATA_FONT

    ws.cell(row=4, column=1, value="PRICING").font = SUBTITLE_FONT
    headers = ["Tier", "Managed Nodes", "Standard ($/yr)", "Professional ($/yr)"]
    for i, h in enumerate(headers, 1):
        ws.cell(row=5, column=i, value=h)
    style_header_row(ws, 5, 4)
    pricing = [
        ("S", "Up to 500", 50000, 100000),
        ("M", "501 – 2,000", 90000, 175000),
        ("L", "2,001 – 5,000", 150000, 250000),
        ("XL", "5,001+", 200000, 350000),
    ]
    for ri, (tier, nodes, std, pro) in enumerate(pricing, 6):
        ws.cell(row=ri, column=1, value=tier).font = DATA_FONT
        ws.cell(row=ri, column=2, value=nodes).font = DATA_FONT
        style_data_cell(ws, ri, 3, MONEY_FMT).value = std
        style_data_cell(ws, ri, 4, MONEY_FMT).value = pro

    ws.cell(row=11, column=1, value="KEY ASSUMPTIONS").font = SUBTITLE_FONT
    assumptions = [
        ("Std → Pro Upgrade Rate", "0% Y1, 15% Y2, 25% Y3, 35% Y4-Y5"),
        ("Annual Logo Churn", "5% Y1, 8% Y2-Y3, 10% Y4-Y5"),
        ("Annual Price Increase", "0% Y1-Y2, 3% Y3-Y4, 5% Y5"),
        ("XL Professional Blended", "$350,000/yr"),
        ("Managed Node Definition", "Agent-deployed endpoints only"),
        ("Free Discovery", "Network devices, desktops, agentless"),
    ]
    for ri, (label, value) in enumerate(assumptions, 12):
        ws.cell(row=ri, column=1, value=label).font = Font(name="Calibri", size=10, bold=True)
        ws.cell(row=ri, column=2, value=value).font = DATA_FONT

    ws.cell(row=19, column=1, value="PACKAGE CONTENTS").font = SUBTITLE_FONT
    ws.cell(row=20, column=1, value="Standard").font = Font(name="Calibri", size=10, bold=True)
    ws.cell(row=20, column=2, value="Pathfinder Discovery Engine + CMDB Ops").font = DATA_FONT
    ws.cell(row=21, column=1, value="Professional").font = Font(name="Calibri", size=10, bold=True)
    ws.cell(row=21, column=2, value="Standard + Integration Intelligence + Service Map Intelligence").font = DATA_FONT

    for col in range(1, 5):
        ws.column_dimensions[get_column_letter(col)].width = 25

    # ── Sheet 2: Likely (Base) ──
    ws2 = wb.create_sheet("Likely (Base)")
    ws2.cell(row=1, column=1, value="Likely Case — Penetration Pricing (Base Plan)").font = TITLE_FONT

    cols = ["Metric", "Year 1", "Year 2", "Year 3", "Year 4", "Year 5"]
    for i, c in enumerate(cols, 1):
        ws2.cell(row=3, column=i, value=c)
    style_header_row(ws2, 3, 6)

    data = [
        ("New Customers", 39, 60, 96, 138, 213),
        ("Total Customers", 39, 94, 180, 301, 546),
        ("Std→Pro Upgrade %", 0, 0.15, 0.25, 0.35, 0.35),
        ("Logo Churn %", 0.05, 0.08, 0.08, 0.10, 0.10),
        ("Price Increase %", 0, 0, 0.03, 0.03, 0.05),
        ("End ARR ($K)", 3000, 10200, 24500, 48000, 84000),
        ("Revenue ($K)", 1500, 6600, 17000, 36000, 66000),
        ("OpEx ($K)", 2200, 3200, 4200, 5800, 6500),
        ("Cash Flow ($K)", -700, 3400, 12800, 30200, 59500),
        ("EBITDA Margin", None, None, None, None, 0.916),
        ("Headcount", 7, 10, 14, 16, 16),
        ("Compass Partners", 0, 20, 45, 70, 100),
    ]
    for ri, row_data in enumerate(data, 4):
        ws2.cell(row=ri, column=1, value=row_data[0]).font = Font(name="Calibri", size=10, bold=True)
        ws2.cell(row=ri, column=1).border = thin_border
        for ci, val in enumerate(row_data[1:], 2):
            cell = style_data_cell(ws2, ri, ci)
            cell.value = val
            if val is not None:
                if "%" in row_data[0] or "Margin" in row_data[0]:
                    cell.number_format = PCT_FMT
                elif "$K" in row_data[0]:
                    cell.number_format = MONEY_FMT

    for col in range(1, 7):
        ws2.column_dimensions[get_column_letter(col)].width = 18

    # ── Sheet 3: Scenario Comparison ──
    ws3 = wb.create_sheet("Scenario Comparison")
    ws3.cell(row=1, column=1, value="Four-Scenario Analysis").font = TITLE_FONT

    cols3 = ["Metric", "Bear", "Likely (Base)", "Bull", "Best Case"]
    for i, c in enumerate(cols3, 1):
        ws3.cell(row=3, column=i, value=c)
    style_header_row(ws3, 3, 5)

    scenarios = [
        ("Y5 Customers", 217, 546, 710, 875),
        ("Y5 ARR ($K)", 33000, 84000, 110000, 135000),
        ("Y5 Revenue ($K)", 25000, 66000, 86000, 106000),
        ("Y5 EBITDA Margin", 0.896, 0.916, 0.923, 0.928),
        ("Profitable in Y1?", "No", "Yes", "Yes", "Yes"),
        ("Std→Pro Y5 Upgrade", 0.25, 0.35, 0.40, 0.45),
        ("Y5 Logo Churn", 0.12, 0.10, 0.08, 0.06),
        ("Y5 Headcount", 14, 16, 18, 20),
        ("Y5 ARR/Employee ($K)", 2357, 5250, 6111, 6750),
    ]
    for ri, row_data in enumerate(scenarios, 4):
        ws3.cell(row=ri, column=1, value=row_data[0]).font = Font(name="Calibri", size=10, bold=True)
        ws3.cell(row=ri, column=1).border = thin_border
        for ci, val in enumerate(row_data[1:], 2):
            cell = style_data_cell(ws3, ri, ci)
            cell.value = val
            if val is not None and isinstance(val, (int, float)):
                if "%" in row_data[0] or "Margin" in row_data[0] or "Churn" in row_data[0] or "Upgrade" in row_data[0]:
                    cell.number_format = PCT_FMT
                elif "$K" in row_data[0]:
                    cell.number_format = MONEY_FMT

    for col in range(1, 6):
        ws3.column_dimensions[get_column_letter(col)].width = 22

    # ── Sheet 4: Cost Model ──
    ws4 = wb.create_sheet("Cost Model")
    ws4.cell(row=1, column=1, value="Cost Model (Likely Case)").font = TITLE_FONT

    cost_cols = ["Category", "Year 1", "Year 2", "Year 3", "Year 4", "Year 5"]
    for i, c in enumerate(cost_cols, 1):
        ws4.cell(row=3, column=i, value=c)
    style_header_row(ws4, 3, 6)

    costs = [
        ("New Hires", 375, 320, 250, 340, 455),
        ("Existing Team", 420, 420, 420, 420, 420),
        ("Prior Salaries", 0, 375, 695, 945, 1290),
        ("Infrastructure", 48, 60, 84, 108, 144),
        ("AI Tokens (Claude)", 24, 48, 72, 96, 120),
        ("Legal + Patent", 90, 15, 10, 25, 10),
        ("Marketing + Partner", 65, 55, 70, 90, 110),
        ("Self-Serve + Support", 0, 100, 200, 350, 500),
        ("Total OpEx", 2200, 3200, 4200, 5800, 6500),
    ]
    for ri, row_data in enumerate(costs, 4):
        is_total = row_data[0] == "Total OpEx"
        ws4.cell(row=ri, column=1, value=row_data[0]).font = Font(
            name="Calibri", size=10, bold=is_total
        )
        ws4.cell(row=ri, column=1).border = thin_border
        for ci, val in enumerate(row_data[1:], 2):
            cell = style_data_cell(ws4, ri, ci, MONEY_FMT)
            cell.value = val
            if is_total:
                cell.font = Font(name="Calibri", size=10, bold=True)

    for col in range(1, 7):
        ws4.column_dimensions[get_column_letter(col)].width = 20

    xlsx_path = os.path.join(ICLOUD_FINANCIAL, "Pathfinder_Financial_Model.xlsx")
    wb.save(xlsx_path)
    size_kb = os.path.getsize(xlsx_path) // 1024
    print(f"Generated XLSX: {xlsx_path} ({size_kb} KB)\n")


# ════════════════════════════════════════════
# PPTX Summary Presentation
# ════════════════════════════════════════════

def generate_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    LIME = RGBColor(0xC8, 0xFF, 0x00)
    DARK = RGBColor(0x0C, 0x0C, 0x18)
    WHITE = RGBColor(0xF0, 0xF0, 0xF8)
    GRAY = RGBColor(0x6A, 0x6A, 0x84)
    TEAL = RGBColor(0x00, 0xD4, 0xAA)

    def add_slide(title_text, subtitle_text=None):
        layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(layout)
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = DARK

        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE

        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.size = Pt(18)
            p2.font.color.rgb = GRAY

        return slide

    def add_text(slide, left, top, width, height, text, size=14, color=WHITE, bold=False):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        return tf

    def add_kpi(slide, left, top, value, label, color=LIME):
        add_text(slide, left, top, 2.5, 0.5, value, size=28, color=color, bold=True)
        add_text(slide, left, top + 0.5, 2.5, 0.3, label, size=12, color=GRAY)

    # ── Slide 1: Title ──
    s1 = add_slide("Avennorth Pathfinder", "Five-Year Business Case  |  March 2026")
    add_text(s1, 0.8, 2.5, 10, 1,
             "CMDB-first integration discovery + AI-powered intelligence for ServiceNow",
             size=20, color=GRAY)
    add_text(s1, 0.8, 3.5, 10, 0.5,
             "Two packages: Standard ($50K+) and Professional ($100K+)  |  Annual by Managed Node count",
             size=16, color=TEAL)

    # ── Slide 2: Pricing ──
    s2 = add_slide("Pricing Structure", "Annual subscription by Managed Node tier")
    pricing_text = (
        "S tier (up to 500 nodes):     Standard $50K/yr    Professional $100K/yr\n"
        "M tier (501 - 2,000):            Standard $90K/yr    Professional $175K/yr\n"
        "L tier (2,001 - 5,000):          Standard $150K/yr  Professional $250K/yr\n"
        "XL tier (5,001+):                  Standard $200K+     Professional Custom"
    )
    add_text(s2, 0.8, 2.0, 11, 3, pricing_text, size=18, color=WHITE)
    add_text(s2, 0.8, 5.0, 11, 1,
             "Managed Node = any endpoint running a Pathfinder agent. "
             "Network devices, desktops, and agentless discovery are FREE.",
             size=14, color=TEAL, bold=True)

    # ── Slide 3: Packages ──
    s3 = add_slide("Two Packages")
    add_text(s3, 0.8, 1.8, 5, 0.4, "STANDARD", size=24, color=LIME, bold=True)
    add_text(s3, 0.8, 2.4, 5, 2,
             "Pathfinder Discovery Engine\n"
             "  - eBPF / ETW agent deployment\n"
             "  - Behavioral confidence scoring\n\n"
             "CMDB Ops\n"
             "  - Automated hygiene\n"
             "  - CI lifecycle management\n"
             "  - Stale record detection & cleanup\n\n"
             "Agentless network discovery (free)",
             size=14, color=WHITE)

    add_text(s3, 7, 1.8, 5, 0.4, "PROFESSIONAL", size=24, color=TEAL, bold=True)
    add_text(s3, 7, 2.4, 5, 2,
             "Everything in Standard, plus:\n\n"
             "Integration Intelligence\n"
             "  - Cross-platform data flow analysis\n"
             "  - Integration anomaly detection\n\n"
             "Service Map Intelligence\n"
             "  - Dependency mapping & validation\n"
             "  - Unmapped service detection",
             size=14, color=WHITE)

    # ── Slide 4: Likely Case ──
    s4 = add_slide("Likely Case (Base Plan)", "Penetration pricing: $50K entry removes procurement friction")
    add_kpi(s4, 0.8, 2.0, "~$84M", "Y5 ARR")
    add_kpi(s4, 3.5, 2.0, "546", "Y5 Customers", TEAL)
    add_kpi(s4, 6.2, 2.0, "91.6%", "Y5 EBITDA Margin", TEAL)
    add_kpi(s4, 8.9, 2.0, "~16", "Y5 Headcount", LIME)
    add_kpi(s4, 11.2, 2.0, "$5.25M", "ARR/Employee", LIME)

    add_text(s4, 0.8, 3.5, 11, 2,
             "Year 1: 39 customers, $3M ARR — profitable in Y1\n"
             "Year 2: 94 customers, $10.2M ARR — 15% Std→Pro upgrade\n"
             "Year 3: 180 customers, $24.5M ARR — 25% upgrade, 3% price increase\n"
             "Year 4: 301 customers, $48M ARR — 35% upgrade, portfolio flywheel\n"
             "Year 5: 546 customers, $84M ARR — 91.6% EBITDA margin",
             size=16, color=WHITE)

    # ── Slide 5: Scenarios ──
    s5 = add_slide("Scenario Analysis", "Bear / Likely / Bull / Best Case")
    scenarios = [
        ("Bear", "~$33M", "217", "89.6%", "No"),
        ("Likely (Base)", "~$84M", "546", "91.6%", "Yes"),
        ("Bull", "~$110M", "710", "92.3%", "Yes"),
        ("Best Case", "~$135M", "875", "92.8%", "Yes"),
    ]
    labels = ["Y5 ARR", "Customers", "EBITDA", "Y1 Profit?"]
    y_start = 2.2
    for i, label in enumerate(labels):
        add_text(s5, 0.8, y_start + i * 0.7, 2, 0.4, label, size=14, color=GRAY, bold=True)
    for j, (name, arr, cust, ebitda, y1) in enumerate(scenarios):
        x = 3 + j * 2.7
        add_text(s5, x, 1.6, 2.5, 0.4, name, size=16,
                 color=LIME if name == "Likely (Base)" else WHITE, bold=True)
        add_text(s5, x, y_start, 2.5, 0.4, arr, size=14, color=LIME if name == "Likely (Base)" else WHITE)
        add_text(s5, x, y_start + 0.7, 2.5, 0.4, cust, size=14, color=WHITE)
        add_text(s5, x, y_start + 1.4, 2.5, 0.4, ebitda, size=14, color=WHITE)
        add_text(s5, x, y_start + 2.1, 2.5, 0.4, y1, size=14,
                 color=TEAL if y1 == "Yes" else RGBColor(0xFF, 0x40, 0x60))

    # ── Slide 6: Strategic Context ──
    s6 = add_slide("Strategic Context")
    add_text(s6, 0.8, 1.8, 11, 0.4, "COREX RELATIONSHIP", size=18, color=LIME, bold=True)
    add_text(s6, 0.8, 2.4, 11, 1,
             "Pathfinder does NOT cannibalize CoreX — it makes CoreX strategically indispensable.\n"
             "CoreX consultants armed with Pathfinder become 'the team that actually fixed the CMDB.'",
             size=14, color=WHITE)

    add_text(s6, 0.8, 3.8, 11, 0.4, "PORTFOLIO FLYWHEEL", size=18, color=TEAL, bold=True)
    add_text(s6, 0.8, 4.4, 11, 1,
             "Pathfinder → surfaces CMDB debt → sells Bearing\n"
             "Pathfinder + Bearing → reveals service map gaps → sells Contour\n"
             "Clean CMDB + mapped services → enables AI incident mgmt → sells Vantage",
             size=14, color=WHITE)

    add_text(s6, 0.8, 5.8, 11, 0.4, "UPGRADE TRIGGER DESIGN", size=18, color=LIME, bold=True)
    add_text(s6, 0.8, 6.4, 11, 1,
             "Standard dashboards show 'detected but not actionable without Professional' items.\n"
             "CoreX consultants point to these dashboards during quarterly reviews — upgrade writes itself.",
             size=14, color=WHITE)

    pptx_path = os.path.join(ICLOUD_PRESENTATIONS, "Pathfinder_Business_Case.pptx")
    prs.save(pptx_path)
    size_kb = os.path.getsize(pptx_path) // 1024
    print(f"Generated PPTX: {pptx_path} ({size_kb} KB)\n")


# ════════════════════════════════════════════
# PDF Generation (via existing scripts)
# ════════════════════════════════════════════

def generate_pdfs():
    from fpdf import FPDF

    class AvennorthPDF(FPDF):
        def header(self):
            self.set_font("Helvetica", "B", 8)
            self.set_text_color(100, 100, 120)
            self.cell(0, 6, "AVENNORTH  |  CONFIDENTIAL", align="L")
            self.ln(8)

        def footer(self):
            self.set_y(-15)
            self.set_font("Helvetica", "", 7)
            self.set_text_color(100, 100, 120)
            self.cell(0, 10, "Avennorth Pathfinder  |  Page %d" % self.page_no(), align="C")

        def h1(self, text):
            self.set_font("Helvetica", "B", 18)
            self.set_text_color(20, 20, 30)
            self.cell(0, 12, text, new_x="LMARGIN", new_y="NEXT")
            self.ln(2)

        def h2(self, text, color=(0, 150, 80)):
            self.set_font("Helvetica", "B", 13)
            self.set_text_color(*color)
            self.cell(0, 9, text, new_x="LMARGIN", new_y="NEXT")
            self.set_draw_color(*color)
            self.line(self.get_x(), self.get_y(), self.get_x() + 50, self.get_y())
            self.ln(4)

        def p(self, text):
            self.set_font("Helvetica", "", 9)
            self.set_text_color(40, 40, 50)
            self.multi_cell(0, 5, text)
            self.ln(2)

        def row(self, label, value, note=""):
            self.set_font("Helvetica", "", 9)
            self.set_text_color(80, 80, 100)
            self.cell(55, 6, label)
            self.set_font("Helvetica", "B", 10)
            self.set_text_color(20, 20, 30)
            self.cell(35, 6, str(value))
            if note:
                self.set_font("Helvetica", "", 8)
                self.set_text_color(120, 120, 140)
                self.cell(0, 6, note)
            self.ln(6)

        def th(self, cols, widths):
            self.set_font("Helvetica", "B", 7)
            self.set_text_color(255, 255, 255)
            self.set_fill_color(30, 30, 50)
            for i, col in enumerate(cols):
                self.cell(widths[i], 6, col, border=1, fill=True, align="C")
            self.ln()

        def tr(self, cells, widths):
            self.set_font("Helvetica", "", 7)
            self.set_text_color(40, 40, 50)
            for i, cell in enumerate(cells):
                self.cell(widths[i], 5, str(cell), border=1, align="C" if i > 0 else "L")
            self.ln()

    # ── Rate Card PDF ──
    pdf = AvennorthPDF()
    pdf.set_auto_page_break(auto=True, margin=20)
    pdf.add_page()

    pdf.h1("Avennorth Pathfinder")
    pdf.p("Pricing & Packaging  |  Effective 2026-03-31")
    pdf.ln(4)

    pdf.h2("Packages", (0, 180, 140))
    w = [35, 80, 75]
    pdf.th(["Package", "Contents", "Key Capabilities"], w)
    for r in [
        ("Standard", "Discovery Engine + CMDB Ops", "eBPF/ETW agents, confidence scoring, CI lifecycle, stale cleanup"),
        ("Professional", "Standard + Integ Intel + Svc Map Intel", "Data flow analysis, anomaly detection, dependency mapping"),
    ]:
        pdf.tr(r, w)
    pdf.ln(4)

    pdf.h2("Pricing by Managed Node Tier", (200, 200, 0))
    w2 = [20, 45, 55, 55]
    pdf.th(["Tier", "Managed Nodes", "Standard (starting at)", "Professional (starting at)"], w2)
    for r in [
        ("S", "Up to 500", "$50,000/yr", "$100,000/yr"),
        ("M", "501 - 2,000", "$90,000/yr", "$175,000/yr"),
        ("L", "2,001 - 5,000", "$150,000/yr", "$250,000/yr"),
        ("XL", "5,001+", "Starting at $200,000/yr", "Custom"),
    ]:
        pdf.tr(r, w2)
    pdf.ln(4)

    pdf.h2("Managed Node Definition", (0, 200, 200))
    pdf.p(
        "A Managed Node is any endpoint where the Pathfinder agent is deployed: "
        "physical servers, VMs, cloud compute instances (EC2, Azure VMs, GCE), "
        "or any endpoint running a supported OS (Linux via eBPF, Windows via ETW)."
    )
    pdf.p(
        "NOT counted (included free): network devices, endpoints/desktops, "
        "any infrastructure discovered through agentless observation."
    )
    pdf.p('Sales one-liner: "If our agent runs on it, it\'s a Managed Node. Everything else we discover is free."')

    pdf.ln(4)
    pdf.h2("Scenario Summary", (60, 200, 120))
    w3 = [40, 30, 30, 30, 30]
    pdf.th(["Metric", "Bear", "Likely (Base)", "Bull", "Best Case"], w3)
    for r in [
        ("Y5 ARR", "~$33M", "~$84M", "~$110M", "~$135M"),
        ("Y5 Customers", "217", "546", "710", "875"),
        ("Y5 EBITDA", "89.6%", "91.6%", "92.3%", "92.8%"),
    ]:
        pdf.tr(r, w3)

    rate_card_path = os.path.join(ICLOUD_DOCS, "Pathfinder_Rate_Card.pdf")
    pdf.output(rate_card_path)
    size_kb = os.path.getsize(rate_card_path) // 1024
    print(f"Generated PDF: {rate_card_path} ({size_kb} KB)")

    # Also run the existing PDF generators to update in-repo PDFs
    import subprocess
    for script in [
        os.path.join(DOCS_DIR, "business-models", "generate_pdfs.py"),
        os.path.join(DOCS_DIR, "guides", "generate_guide_pdfs.py"),
    ]:
        if os.path.exists(script):
            result = subprocess.run(
                ["python3", script],
                capture_output=True, text=True, cwd=PROJECT_ROOT
            )
            if result.returncode == 0:
                print(f"  Ran: {os.path.basename(script)} — {result.stdout.strip()}")
            else:
                print(f"  WARN: {os.path.basename(script)} — {result.stderr.strip()}")
    print()


# ════════════════════════════════════════════
# Main
# ════════════════════════════════════════════

if __name__ == "__main__":
    print("=" * 60)
    print("Avennorth Pathfinder — Generating iCloud Output Files")
    print("=" * 60)
    print()

    generate_docx_files()
    generate_xlsx()
    generate_pptx()
    generate_pdfs()

    print("=" * 60)
    print("All output files generated.")
    print(f"  Docs:          {ICLOUD_DOCS}")
    print(f"  Financial:     {ICLOUD_FINANCIAL}")
    print(f"  Presentations: {ICLOUD_PRESENTATIONS}")
    print("=" * 60)
